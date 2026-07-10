from flask import Flask, render_template, jsonify, request, send_file
import pandas as pd
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import os
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
DATA_PATH = "dataset/final/boutique_intelligence.csv"
DATA_CACHE = None

# LOAD DATA and AI ENGINE

def load_data():
    global DATA_CACHE
    if DATA_CACHE is not None:
        return DATA_CACHE.copy()
    df = pd.read_csv(DATA_PATH)
    numeric_cols = [
        "Average Rating",
        "Review Count",
        "Digital_Score",
        "Opportunity_Score",
        "Has_Website",
        "Has_Instagram",
        "Latitude",
        "Longitude"
    ]
    for col in numeric_cols:
        if col in df.columns:
            df[col] = pd.to_numeric(
                df[col],
                errors="coerce"
            ).fillna(0)


    df = df.fillna("")


    # Growth Priority Model

    def growth_priority(row):

        score = row["Opportunity_Score"]

        if score >= 80:
            return "🔥 Critical Growth Lead"

        elif score >= 60:
            return "🚀 High Potential"

        elif score >= 35:
            return "⚡ Medium Opportunity"

        else:
            return "🟢 Stable Business"


    df["Growth_Priority"] = df.apply(
        growth_priority,
        axis=1
    )


    # Digital Risk Engine

    def digital_risk(row):

        risk = 0

        if row["Has_Website"] == 0:
            risk += 40

        if row["Has_Instagram"] == 0:
            risk += 30

        if row["Digital_Score"] < 50:
            risk += 30


        if risk >= 70:
            return "High Risk"

        elif risk >= 40:
            return "Medium Risk"

        else:
            return "Low Risk"


    df["Digital_Risk"] = df.apply(
        digital_risk,
        axis=1
    )


    # AI Recommendation Engine

    def ai_strategy(row):

        plan = []


        if row["Has_Website"] == 0:

            plan.append(
                "Launch responsive website"
            )


        if row["Has_Instagram"] == 0:

            plan.append(
                "Build Instagram presence"
            )


        if row["Digital_Score"] < 60:

            plan.append(
                "Improve SEO and Google visibility"
            )


        if row["Review Count"] < 50:

            plan.append(
                "Increase customer reviews"
            )


        if not plan:

            plan.append(
                "Maintain digital growth strategy"
            )


        return " | ".join(plan)


    df["AI_Recommendation"] = df.apply(
        ai_strategy,
        axis=1
    )


    # AI Roadmap

    df["Growth_Roadmap"] = (
        "Week 1: Digital Audit | "
        "Week 2: Website & SEO | "
        "Week 3: Social Campaign | "
        "Week 4: Lead Generation"
    )


    DATA_CACHE = df.copy()

    return df



# ===============================
# DASHBOARD PAGE
# ===============================

@app.route("/")
def dashboard():

    df = load_data()

    total = len(df)

    avg_rating = round(
        df["Average Rating"].mean(),
        2
    )


    no_website = round(
        (df["Has_Website"] == 0)
        .mean() * 100,
        1
    )


    hidden = int(
        df["Business_Segment"]
        .eq("Hidden Gem")
        .sum()
    )


    social_presence = round(
        df["Has_Instagram"]
        .mean() * 100,
        1
    )


    opportunities = int(
        (df["Opportunity_Score"] >= 75)
        .sum()
    )


    website_count = (
        df["Has_Website"]
        .replace({
            0:"No Website",
            1:"Has Website"
        })
        .value_counts()
    )


    website_data = {
        "labels":list(website_count.index),
        "values":[
            int(x)
            for x in website_count.values
        ]
    }


    segment_count = (
        df["Business_Segment"]
        .value_counts()
    )


    segment_data = {
        "labels":list(segment_count.index),
        "values":[
            int(x)
            for x in segment_count.values
        ]
    }


    return render_template(
        "index.html",
        total=total,
        avg_rating=avg_rating,
        no_website=no_website,
        hidden=hidden,
        social_presence=social_presence,
        opportunities=opportunities,
        website_data=website_data,
        segment_data=segment_data
    )
# ===============================
# BUSINESS DATA API
# ===============================

@app.route("/api/businesses")
def businesses():

    df = load_data()

    return jsonify(
        df.to_dict(
            orient="records"
        )
    )



# ===============================
# SMART FILTER API
# ===============================

@app.route("/api/filter")
def filter_business():

    df = load_data()


    segment = request.args.get("segment")
    website = request.args.get("website")
    priority = request.args.get("priority")


    if segment:

        df = df[
            df["Business_Segment"] == segment
        ]


    if website == "missing":

        df = df[
            df["Has_Website"] == 0
        ]


    if website == "available":

        df = df[
            df["Has_Website"] == 1
        ]


    if priority:

        df = df[
            df["Growth_Priority"]
            .str.contains(priority)
        ]


    return jsonify(
        df.to_dict(
            orient="records"
        )
    )



# ===============================
# LOCATION INTELLIGENCE
# ===============================
@app.route("/api/location-summary")
def location_summary():

    df = load_data()

    zones = df["Municipality"].nunique()

    top_area = (
        df.groupby("Municipality")
        ["Opportunity_Score"]
        .mean()
        .idxmax()
    )

    return jsonify({

        "zones": int(zones),

        "top_area": top_area

    })
@app.route("/api/locations")
def locations():

    df = load_data()


    location_data = (

        df.groupby("Municipality")

        .agg(

            Businesses=(
                "Name",
                "count"
            ),

            Avg_Rating=(
                "Average Rating",
                "mean"
            ),

            Avg_Opportunity=(
                "Opportunity_Score",
                "mean"
            ),

            Avg_Digital=(
                "Digital_Score",
                "mean"
            )

        )

        .reset_index()

    )


    location_data = (

        location_data
        .round(2)

    )


    return jsonify(

        location_data
        .to_dict(
            orient="records"
        )

    )



# ===============================
# MAP DATA API
# ===============================

@app.route("/api/map")
def map_data():

    df = load_data()


    columns = [

        "Name",

        "Latitude",

        "Longitude",

        "Average Rating",

        "Opportunity_Score",

        "Digital_Score",

        "Growth_Priority"

    ]


    return jsonify(

        df[columns]

        .to_dict(
            orient="records"
        )

    )




# ===============================
# SCATTER ANALYSIS API
# ===============================

@app.route("/api/scatter")
def scatter():


    df = load_data()


    result = []


    for _, row in df.iterrows():


        result.append({

            "x":

            float(
                row["Average Rating"]
            ),


            "y":

            float(
                row["Opportunity_Score"]
            )

        })


    return jsonify(result)






# ===============================
# EXECUTIVE INSIGHTS API
# ===============================

@app.route("/api/insights")
def insights():

    df = load_data()


    total = len(df)


    missing_site = int(

        (
            df["Has_Website"] == 0
        )

        .sum()

    )


    high_growth = int(

        (
            df["Opportunity_Score"] >= 75
        )

        .sum()

    )


    top = (

        df.sort_values(

            "Opportunity_Score",

            ascending=False

        )

        .iloc[0]

    )



    data = [


        {

            "title":

            "🔥 Highest Opportunity",


            "value":

            top["Name"],


            "text":

            f"""
            Opportunity Score:
            {top['Opportunity_Score']}%
            """

        },


        {

            "title":

            "🌐 Digital Gap",


            "value":

            f"{missing_site} Businesses",


            "text":

            "Businesses require websites and online visibility improvement"

        },



        {

            "title":

            "🚀 Growth Market",


            "value":

            f"{high_growth} Leads",


            "text":

            "AI detected businesses ready for digital transformation"

        },


        {

            "title":

            "🤖 AI Strategy",


            "value":

            "Recommended",


            "text":

            "Website + SEO + Social Media + Customer Acquisition"

        }

    ]


    return jsonify(data)

# =====================================
# EXECUTIVE SUMMARY API
# =====================================

@app.route("/api/executive-summary")
def executive_summary():

    df = load_data()

    total = len(df)

    no_website = int(
        (df["Has_Website"] == 0).sum()
    )

    website_gap = round(
        (no_website / total) * 100,
        1
    )

    hidden = int(
        (df["Business_Segment"] == "Hidden Gem")
        .sum()
    )

    top_area = (
        df.groupby("Municipality")
        ["Opportunity_Score"]
        .mean()
        .idxmax()
    )

    avg_rating = round(
        df["Average Rating"].mean(),
        2
    )

    avg_digital = round(
        df["Digital_Score"].mean(),
        1
    )

    summary = f"""
Threads To Trends analyzed {total} fashion businesses across multiple locations.

{website_gap}% of businesses currently operate without a website, creating a significant digital growth opportunity.

The average customer rating is {avg_rating}, indicating strong customer trust despite limited digital adoption.

AI identified {hidden} Hidden Gem businesses requiring immediate digital transformation.

{top_area} emerged as the highest opportunity market for website development, SEO and social media growth.

Overall recommendation:
Prioritize businesses with high ratings but weak digital presence to maximize business growth.
"""

    return jsonify({
        "summary": summary.strip()
    })




# ===============================
# COMPETITOR GAP ANALYSIS
# ===============================

@app.route("/api/competitor/<int:id>")
def competitor(id):


    df = load_data()


    business = df.iloc[id]


    market_average = round(

        df["Digital_Score"]

        .mean(),

        2

    )


    gap = round(

        market_average

        -

        business["Digital_Score"],

        2

    )


    return jsonify({


        "business":

        business["Name"],


        "current_score":

        float(
            business["Digital_Score"]
        ),


        "market_average":

        market_average,


        "gap":

        gap,


        "risk":

        business["Digital_Risk"],


        "recommendation":

        business["AI_Recommendation"]


    })
# ===============================
# PREMIUM PDF AUDIT REPORT
# ===============================

def clean_pdf_text(text):

    text = str(text)

    remove = [
        "🔥","🚀","💎","📊","📍",
        "⭐","🌐","📸","🤖",
        "⚠️","🟢","⚡"
    ]

    for x in remove:
        text = text.replace(x,"")

    return text

@app.route("/api/pdf/<int:id>")
def pdf_report(id):

    df = load_data()

    b = df.iloc[id]

    pdf = FPDF()

    pdf.add_page()


    # HEADER

    pdf.set_fill_color(
        185,131,255
    )

    pdf.rect(
        0,
        0,
        210,
        35,
        "F"
    )


    pdf.set_text_color(
        255,255,255
    )


    pdf.set_font(
        "Arial",
        "B",
        22
    )


    pdf.cell(
        0,
        15,
        "THREADS TO TRENDS",
        ln=True,
        align="C"
    )


    pdf.set_font(
        "Arial",
        "",
        12
    )


    pdf.cell(
        0,
        10,
        "AI Business Intelligence Audit Report",
        ln=True,
        align="C"
    )


    pdf.ln(20)


    pdf.set_text_color(
        0,0,0
    )


    # BUSINESS INFO

    pdf.set_font(
        "Arial",
        "B",
        16
    )


    pdf.cell(
        0,
        10,
        "Business Profile",
        ln=True
    )


    pdf.set_font(
        "Arial",
        "",
        12
    )


    info=[

        f"Name : {b['Name']}",

        f"Location : {b['Municipality']}",

        f"Rating : {b['Average Rating']}",

        f"Digital Score : {b['Digital_Score']}/100",

        f"Opportunity Score : {b['Opportunity_Score']}%",

        f"Priority : {b['Growth_Priority']}"

    ]


    for item in info:

        pdf.multi_cell(
            0,
            8,
            clean_pdf_text(item)
        )


    pdf.ln(5)


    # SCORE BAR

    pdf.set_font(
        "Arial",
        "B",
        14
    )

    pdf.cell(
        0,
        10,
        "Digital Growth Analysis",
        ln=True
    )


    pdf.set_fill_color(
        185,131,255
    )


    pdf.cell(
        int(b["Opportunity_Score"]),
        8,
        "",
        fill=True
    )


    pdf.ln(15)



    # AI SECTION

    pdf.set_font(
        "Arial",
        "B",
        14
    )


    pdf.cell(
        0,
        10,
        "AI Recommendation",
        ln=True
    )


    pdf.set_font(
        "Arial",
        "",
        12
    )


    pdf.multi_cell(
        0,
        8,
        clean_pdf_text(
            b["AI_Recommendation"]
        )
    )


    pdf.ln(5)


    pdf.set_font(
        "Arial",
        "B",
        14
    )


    pdf.cell(
        0,
        10,
        "30 Day Growth Roadmap",
        ln=True
    )


    pdf.set_font(
        "Arial",
        "",
        12
    )


    pdf.multi_cell(
        0,
        8,
        clean_pdf_text(
            b["Growth_Roadmap"]
        )
    )



    pdf.ln(15)


    pdf.set_font(
        "Arial",
        "I",
        10
    )


    pdf.cell(
        0,
        10,
        "Generated using Threads To Trends AI Intelligence Engine",
        align="C"
    )


    path="Threads_AI_Audit.pdf"

    pdf.output(path)


    return send_file(
        path,
        as_attachment=True
    )

# ===============================
# EXECUTIVE PPT ENGINE V2
# ===============================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ---------- Theme Colors ----------

BG = RGBColor(20,15,30)

CARD = RGBColor(58,40,88)

ACCENT = RGBColor(255,209,102)

WHITE = RGBColor(255,255,255)

SUB = RGBColor(210,210,210)


# ---------- Background ----------

def set_background(slide):

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


# ---------- Title ----------

def add_title(slide, text):

    box = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.25),
        Inches(12),
        Inches(0.6)
    )

    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]

    p.text = text

    p.alignment = PP_ALIGN.LEFT

    p.font.name = "Calibri"

    p.font.size = Pt(28)

    p.font.bold = True

    p.font.color.rgb = ACCENT


# ---------- Paragraph ----------

def add_text(slide,text,x,y,w,h,size=16,color=WHITE):

    box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )

    tf = box.text_frame
    tf.clear()

    tf.word_wrap = True

    p = tf.paragraphs[0]

    p.text = str(text)

    p.alignment = PP_ALIGN.LEFT

    p.font.name = "Calibri"

    p.font.size = Pt(size)

    p.font.color.rgb = color


# ---------- KPI Card ----------

def add_card(slide,x,y,title,value):

    shape = slide.shapes.add_shape(

        MSO_SHAPE.ROUNDED_RECTANGLE,

        Inches(x),

        Inches(y),

        Inches(2.4),

        Inches(1.15)

    )

    shape.fill.solid()

    shape.fill.fore_color.rgb = CARD

    shape.line.color.rgb = ACCENT

    tf = shape.text_frame
    tf.clear()

    p = tf.paragraphs[0]

    p.text = str(value)

    p.alignment = PP_ALIGN.CENTER

    p.font.name = "Calibri"

    p.font.size = Pt(24)

    p.font.bold = True

    p.font.color.rgb = ACCENT

# ---------- Divider ----------

def add_divider(slide, y):

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(y),
        Inches(12),
        Inches(0.02)
    )

    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.color.rgb = ACCENT


# =====================================
# PPT ROUTE
# =====================================

def ppt_report():

    df = load_data()

    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =====================================
    # DATA CALCULATIONS
    # =====================================

    total = len(df)

    website_missing = int(
        (df["Has_Website"] == 0).sum()
    )

    website_gap = round(
        website_missing / total * 100,
        1
    )

    instagram_presence = round(
        df["Has_Instagram"].mean() * 100,
        1
    )

    hidden = int(
        (df["Business_Segment"] == "Hidden Gem").sum()
    )

    avg_rating = round(
        df["Average Rating"].mean(),
        2
    )

    avg_digital = round(
        df["Digital_Score"].mean(),
        1
    )

    avg_opportunity = round(
        df["Opportunity_Score"].mean(),
        1
    )

    # =====================================
    # SLIDE 1
    # =====================================

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    set_background(slide)

    add_title(
        slide,
        "THREADS TO TRENDS"
    )

    add_text(
        slide,
        "AI Powered Fashion Business Intelligence",
        0.8,
        1.25,
        8,
        0.5,
        24
    )

    add_divider(
        slide,
        2.0
    )

    add_text(
        slide,
        "Google Maps Business Analysis\n\n"
        "Digital Growth Opportunity Finder\n\n"
        "Executive Consulting Report",
        0.9,
        2.3,
        8,
        2,
        18,
        SUB
    )

    add_text(
        slide,
        "Prepared using Python • Flask • Machine Learning • Business Intelligence",
        0.9,
        5.8,
        10,
        0.4,
        13,
        SUB
    )

    # =====================================
    # SLIDE 2
    # =====================================

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    set_background(slide)

    add_title(
        slide,
        "Executive Market Overview"
    )

    add_card(
        slide,
        0.5,
        1.2,
        "Businesses",
        total
    )

    add_card(
        slide,
        3.2,
        1.2,
        "Website Gap",
        f"{website_gap}%"
    )

    add_card(
        slide,
        5.9,
        1.2,
        "Hidden Gems",
        hidden
    )

    add_card(
        slide,
        8.6,
        1.2,
        "Avg Rating",
        avg_rating
    )

    add_card(
        slide,
        2.0,
        3.1,
        "Digital Score",
        avg_digital
    )

    add_card(
        slide,
        5.4,
        3.1,
        "Opportunity",
        f"{avg_opportunity}%"
    )

    add_divider(
        slide,
        5.0
    )

    add_text(
        slide,
        "Executive Insight",
        0.7,
        5.25,
        3,
        0.3,
        18
    )

    add_text(
        slide,
        "The analysis indicates that a large percentage of fashion boutiques "
        "maintain strong customer ratings despite lacking professional digital "
        "presence. Businesses with high trust and low digital maturity have been "
        "identified as the highest priority opportunities.",
        0.7,
        5.7,
        11.5,
        1.0,
        14,
        SUB
    )

@app.route("/api/ppt")
def ppt_report():

    df = load_data()

    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =====================================
    # DATA CALCULATIONS
    # =====================================

    total = len(df)

    website_missing = int(
        (df["Has_Website"] == 0).sum()
    )

    website_gap = round(
        website_missing / total * 100,
        1
    )

    instagram_presence = round(
        df["Has_Instagram"].mean() * 100,
        1
    )

    hidden = int(
        (df["Business_Segment"] == "Hidden Gem").sum()
    )

    avg_rating = round(
        df["Average Rating"].mean(),
        2
    )

    avg_digital = round(
        df["Digital_Score"].mean(),
        1
    )

    avg_opportunity = round(
        df["Opportunity_Score"].mean(),
        1
    )

    # =====================================
    # SLIDE 1
    # =====================================

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    set_background(slide)

    add_title(
        slide,
        "THREADS TO TRENDS"
    )

    add_text(
        slide,
        "AI Powered Fashion Business Intelligence",
        0.8,
        1.25,
        8,
        0.5,
        24
    )

    add_divider(
        slide,
        2.0
    )

    add_text(
        slide,
        "Google Maps Business Analysis\n\n"
        "Digital Growth Opportunity Finder\n\n"
        "Executive Consulting Report",
        0.9,
        2.3,
        8,
        2,
        18,
        SUB
    )

    add_text(
        slide,
        "Prepared using Python • Flask • Machine Learning • Business Intelligence",
        0.9,
        5.8,
        10,
        0.4,
        13,
        SUB
    )

    # =====================================
    # SLIDE 2
    # =====================================

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    set_background(slide)

    add_title(
        slide,
        "Executive Market Overview"
    )

    add_card(
        slide,
        0.5,
        1.2,
        "Businesses",
        total
    )

    add_card(
        slide,
        3.2,
        1.2,
        "Website Gap",
        f"{website_gap}%"
    )

    add_card(
        slide,
        5.9,
        1.2,
        "Hidden Gems",
        hidden
    )

    add_card(
        slide,
        8.6,
        1.2,
        "Avg Rating",
        avg_rating
    )

    add_card(
        slide,
        2.0,
        3.1,
        "Digital Score",
        avg_digital
    )

    add_card(
        slide,
        5.4,
        3.1,
        "Opportunity",
        f"{avg_opportunity}%"
    )

    add_divider(
        slide,
        5.0
    )

    add_text(
        slide,
        "Executive Insight",
        0.7,
        5.25,
        3,
        0.3,
        18
    )

    add_text(
        slide,
        "The analysis indicates that a large percentage of fashion boutiques "
        "maintain strong customer ratings despite lacking professional digital "
        "presence. Businesses with high trust and low digital maturity have been "
        "identified as the highest priority opportunities.",
        0.7,
        5.7,
        11.5,
        1.0,
        14,
        SUB
    )
        # =====================================
    # SLIDE 6 : BUSINESS SEGMENTATION
    # =====================================

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    set_background(slide)

    add_title(
        slide,
        "Business Segmentation"
    )

    segments = (
        df["Business_Segment"]
        .value_counts()
        .reset_index()
    )

    segments.columns = [
        "Segment",
        "Count"
    ]

    x = 0.7
    y = 1.4

    for _, row in segments.iterrows():

        add_card(
            slide,
            x,
            y,
            row["Segment"],
            row["Count"]
        )

        x += 2.8

        if x > 9:
            x = 0.7
            y += 1.6


    add_text(
        slide,
        "AI grouped businesses according to their digital maturity and market opportunity.",
        0.7,
        5.8,
        11,
        0.5,
        14,
        SUB
    )


    # =====================================
    # SLIDE 7 : AI RECOMMENDATIONS
    # =====================================

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    set_background(slide)

    add_title(
        slide,
        "30 Day AI Growth Roadmap"
    )

    add_card(slide,0.7,1.2,"Week 1","Audit")
    add_card(slide,3.5,1.2,"Week 2","Website")
    add_card(slide,6.3,1.2,"Week 3","SEO")
    add_card(slide,9.1,1.2,"Week 4","Marketing")


    add_text(
        slide,
        "Digital Audit\nCompetitor Analysis",
        0.7,
        2.8,
        2.2,
        1.2,
        14
    )

    add_text(
        slide,
        "Website\nProduct Catalogue",
        3.5,
        2.8,
        2.2,
        1.2,
        14
    )

    add_text(
        slide,
        "SEO\nGoogle Business\nInstagram",
        6.3,
        2.8,
        2.2,
        1.5,
        14
    )

    add_text(
        slide,
        "Lead Generation\nPerformance Tracking",
        9.1,
        2.8,
        2.2,
        1.5,
        14
    )

    add_divider(
        slide,
        5.0
    )

    add_text(
        slide,
        "Executive Recommendation",
        0.7,
        5.2,
        3,
        0.4,
        18
    )

    recommendation = (
        "Businesses with strong customer ratings but missing websites "
        "should be prioritized for digital transformation. "
        "Launching responsive websites, improving SEO, "
        "optimizing Google Business Profiles and strengthening "
        "Instagram presence can significantly increase online visibility "
        "and customer acquisition."
    )

    add_text(
        slide,
        recommendation,
        0.7,
        5.7,
        11.3,
        1.0,
        14,
        SUB
    )


    # =====================================
    # SLIDE 8 : THANK YOU
    # =====================================

    slide = prs.slides.add_slide(
        prs.slide_layouts[6]
    )

    set_background(slide)

    add_title(
        slide,
        "Thank You"
    )

    add_divider(
        slide,
        1.2
    )

    add_text(
        slide,
        "Threads To Trends",
        0.8,
        2.0,
        8,
        0.5,
        30
    )

    add_text(
        slide,
        "AI Powered Fashion Business Intelligence Platform",
        0.8,
        2.8,
        9,
        0.5,
        18,
        SUB
    )

    add_text(
        slide,
        "Prepared using Python • Flask • Machine Learning • Business Intelligence",
        0.8,
        5.8,
        11,
        0.4,
        14,
        SUB
    )


    # =====================================
    # SAVE REPORT
    # =====================================

    path = "Threads_AI_Report.pptx"

    prs.save(path)

    return send_file(
        path,
        as_attachment=True,
        download_name="Threads_To_Trends_AI_Report.pptx"
    )

if __name__ == "__main__":
 app.run(debug=True)

